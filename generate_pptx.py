from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand Colors ──────────────────────────────────────────────
TEAL        = RGBColor(0x00, 0x56, 0x70)
TEAL_DARK   = RGBColor(0x00, 0x3d, 0x50)
TEAL_LIGHT  = RGBColor(0xE6, 0xF2, 0xF5)
RED         = RGBColor(0x9B, 0x1B, 0x30)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0x44, 0x44, 0x44)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY    = RGBColor(0x88, 0x88, 0x88)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ───────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_para(tf, text, font_size=16, bold=False, color=GRAY,
             align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def header_bar(slide, title, subtitle=None):
    """Teal top bar with slide title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), TEAL)
    add_rect(slide, 0, Inches(1.25), Inches(0.06), SLIDE_H - Inches(1.25), RED)
    add_text(slide, title,
             Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.65),
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.82), Inches(12.5), Inches(0.38),
                 font_size=14, color=RGBColor(0xCC, 0xE8, 0xF0))


def footer(slide, text="AI Ideathon 2026  |  Stewart Title Guaranty Company"):
    add_rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), TEAL_DARK)
    add_text(slide, text,
             Inches(0.3), SLIDE_H - Inches(0.36), Inches(12.7), Inches(0.34),
             font_size=10, color=RGBColor(0xAA, 0xCC, 0xD8), align=PP_ALIGN.LEFT)


def bullet_box(slide, items, l, t, w, h, font_size=15, title=None, title_color=TEAL):
    add_rect(slide, l, t, w, h, LIGHT_GRAY)
    offset = Inches(0.12)
    if title:
        add_text(slide, title, l + Inches(0.15), t + Inches(0.1),
                 w - Inches(0.2), Inches(0.35),
                 font_size=13, bold=True, color=title_color)
        offset = Inches(0.48)
    txBox = slide.shapes.add_textbox(l + Inches(0.15), t + offset,
                                     w - Inches(0.3), h - offset - Inches(0.1))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = GRAY
        run.font.name = "Calibri"


def icon_card(slide, emoji, title, body, l, t, w, h, accent=TEAL):
    add_rect(slide, l, t, w, h, WHITE, line_color=RGBColor(0xCC, 0xCC, 0xCC))
    add_rect(slide, l, t, w, Inches(0.07), accent)
    add_text(slide, emoji, l + Inches(0.15), t + Inches(0.12),
             Inches(0.55), Inches(0.55), font_size=26, color=GRAY)
    add_text(slide, title, l + Inches(0.15), t + Inches(0.58),
             w - Inches(0.3), Inches(0.38),
             font_size=13, bold=True, color=accent)
    txBox = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.95),
                                     w - Inches(0.3), h - Inches(1.05))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY
    run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Full teal background
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, TEAL)

# Red accent stripe left
add_rect(slide, 0, 0, Inches(0.55), SLIDE_H, RED)

# White card center
add_rect(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(4.5), WHITE)

# Red top accent on card
add_rect(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.1), RED)

# Company name top
add_text(slide, "STEWART TITLE GUARANTY COMPANY",
         Inches(0.9), Inches(0.45), Inches(11.5), Inches(0.55),
         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "AI Ideathon 2026",
         Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.42),
         font_size=13, color=RGBColor(0xCC, 0xE8, 0xF0), align=PP_ALIGN.CENTER)

# Main title
add_text(slide, "AI-Governed Enterprise\nKnowledge Assistant",
         Inches(1.1), Inches(1.7), Inches(11.1), Inches(2.0),
         font_size=38, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Subtitle
add_text(slide, "Cited  ·  Traceable  ·  Governed  ·  Enterprise-Grade",
         Inches(1.1), Inches(3.65), Inches(11.1), Inches(0.5),
         font_size=17, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# Divider
add_rect(slide, Inches(4.5), Inches(4.22), Inches(4.33), Inches(0.04), RED)

# Presenter & URL
add_text(slide, "Meet P.  |  meet.mvp83@gmail.com",
         Inches(1.1), Inches(4.38), Inches(11.1), Inches(0.38),
         font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "🌐  aiideathon.vercel.app",
         Inches(1.1), Inches(4.78), Inches(11.1), Inches(0.38),
         font_size=13, color=TEAL, align=PP_ALIGN.CENTER, bold=True)

footer(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(slide, "The Problem", "Enterprise knowledge is locked away — and general AI makes it worse")

# Two problem cards
card_t = Inches(1.55)
card_h = Inches(4.45)
card_w = Inches(6.1)

# Card 1
add_rect(slide, Inches(0.22), card_t, card_w, card_h, TEAL_LIGHT)
add_rect(slide, Inches(0.22), card_t, card_w, Inches(0.08), RED)
add_text(slide, "🔍", Inches(0.35), card_t + Inches(0.15), Inches(0.7), Inches(0.7), font_size=30, color=GRAY)
add_text(slide, "Problem 1: Knowledge is Inaccessible",
         Inches(0.35), card_t + Inches(0.78), card_w - Inches(0.25), Inches(0.5),
         font_size=16, bold=True, color=TEAL)

txBox = slide.shapes.add_textbox(Inches(0.35), card_t + Inches(1.3),
                                  card_w - Inches(0.3), Inches(2.9))
txBox.word_wrap = True
tf = txBox.text_frame
tf.word_wrap = True
items1 = [
    "▸  Hundreds of PDFs & Word files spread across drives",
    "▸  Employees don't know what exists or which version is current",
    "▸  Questions routed to SMEs — answers take days",
    "▸  New staff especially impacted — no institutional knowledge",
    "▸  Productivity lost at scale across the organization",
]
for i, item in enumerate(items1):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = item
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY
    run.font.name = "Calibri"

# Card 2
add_rect(slide, Inches(6.88), card_t, card_w, card_h, RGBColor(0xFD, 0xF0, 0xF2))
add_rect(slide, Inches(6.88), card_t, card_w, Inches(0.08), RED)
add_text(slide, "⚠️", Inches(7.0), card_t + Inches(0.15), Inches(0.7), Inches(0.7), font_size=30, color=GRAY)
add_text(slide, "Problem 2: General AI Is Dangerous Here",
         Inches(7.0), card_t + Inches(0.78), card_w - Inches(0.25), Inches(0.5),
         font_size=16, bold=True, color=RED)

txBox2 = slide.shapes.add_textbox(Inches(7.0), card_t + Inches(1.3),
                                   card_w - Inches(0.3), Inches(2.9))
txBox2.word_wrap = True
tf2 = txBox2.text_frame
tf2.word_wrap = True
items2 = [
    "▸  ChatGPT answers from internet — not Stewart's documents",
    "▸  Hallucination: confident answers that are factually wrong",
    "▸  Zero traceability — no audit trail, no citations",
    "▸  No governance — violates enterprise compliance standards",
    "▸  Cannot verify or defend an AI-generated answer",
]
for i, item in enumerate(items2):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = item
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY
    run.font.name = "Calibri"

# Bottom question
add_rect(slide, Inches(0.22), Inches(6.1), Inches(12.9), Inches(0.82), TEAL)
add_text(slide,
         "\"How do we give employees fast, accurate answers — without hallucination and with full traceability?\"",
         Inches(0.45), Inches(6.18), Inches(12.5), Inches(0.65),
         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

footer(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(slide, "Our Solution", "RAG-Based AI — Answers from approved documents only, with every claim cited")

# RAG explanation — 3 steps
step_y = Inches(1.55)
step_h = Inches(2.1)
step_w = Inches(3.85)
gap = Inches(0.28)

steps = [
    ("1", "Employee Asks a Question", "Types or speaks a natural language question in the web app or Slack. No special syntax or training required.", TEAL),
    ("2", "System Searches Approved Docs", "Discovery Engine semantically searches our indexed document library. Finds the top 3 most relevant passages — by meaning, not just keywords.", RGBColor(0x00, 0x7A, 0x99)),
    ("3", "AI Generates a Cited Answer", "Gemini AI reads only the retrieved passages and generates an answer. Every factual claim is tagged [1] [2] with a source. Confidence rated: FULL / PARTIAL / NO SUPPORT.", RED),
]

for i, (num, title, body, color) in enumerate(steps):
    x = Inches(0.22) + i * (step_w + gap)
    add_rect(slide, x, step_y, step_w, step_h, LIGHT_GRAY)
    add_rect(slide, x, step_y, step_w, Inches(0.08), color)

    # Step number circle (as rectangle)
    add_rect(slide, x + Inches(0.15), step_y + Inches(0.15),
             Inches(0.55), Inches(0.55), color)
    add_text(slide, num, x + Inches(0.15), step_y + Inches(0.13),
             Inches(0.55), Inches(0.55),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, title, x + Inches(0.15), step_y + Inches(0.75),
             step_w - Inches(0.3), Inches(0.55),
             font_size=14, bold=True, color=color)

    txBox = slide.shapes.add_textbox(x + Inches(0.15), step_y + Inches(1.3),
                                     step_w - Inches(0.25), Inches(0.75))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY
    run.font.name = "Calibri"

    # Arrow between steps
    if i < 2:
        add_text(slide, "→",
                 x + step_w + Inches(0.04), step_y + Inches(0.7),
                 gap + Inches(0.1), Inches(0.6),
                 font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Key outcome box
add_rect(slide, Inches(0.22), Inches(3.78), Inches(12.89), Inches(1.22), TEAL)
add_text(slide, "The Result:",
         Inches(0.45), Inches(3.88), Inches(2.0), Inches(0.42),
         font_size=14, bold=True, color=WHITE)
add_text(slide,
         "Every answer is cited  ·  Every source is traceable  ·  Every confidence level is rated  ·  No hallucination by design",
         Inches(0.45), Inches(4.3), Inches(12.5), Inches(0.5),
         font_size=15, color=RGBColor(0xCC, 0xE8, 0xF0), align=PP_ALIGN.LEFT)

# Feature pills
features = ["📎 Numbered Citations", "🎯 Confidence Scoring", "💬 Multi-Turn Memory",
            "💡 Follow-Up Suggestions", "👍 Feedback Rating", "📊 Analytics Dashboard", "🤖 Slack Bot"]
pill_y = Inches(5.15)
pill_x = Inches(0.22)
for feat in features:
    pill_w = Inches(1.72)
    add_rect(slide, pill_x, pill_y, pill_w, Inches(0.42), TEAL_LIGHT)
    add_text(slide, feat, pill_x + Inches(0.06), pill_y + Inches(0.05),
             pill_w - Inches(0.1), Inches(0.34),
             font_size=10.5, color=TEAL, bold=True)
    pill_x += pill_w + Inches(0.1)

footer(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Live Demo Intro
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, TEAL)
add_rect(slide, 0, 0, Inches(0.55), SLIDE_H, RED)

add_text(slide, "LIVE DEMO", Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.3),
         font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "aiideathon.vercel.app",
         Inches(0.8), Inches(3.1), Inches(11.7), Inches(0.7),
         font_size=24, color=RGBColor(0xCC, 0xE8, 0xF0), align=PP_ALIGN.CENTER)

demo_steps = [
    "Ask a question  →  Observe citations & confidence badge",
    "Click a follow-up chip  →  See multi-turn memory",
    "Ask an out-of-scope question  →  See NO SUPPORT response",
    "Rate an answer  →  Thumbs up / down feedback",
    "View Analytics dashboard  →  Confidence & satisfaction breakdown",
    "Slack  →  /ask command live demo",
]
for i, step in enumerate(demo_steps):
    add_text(slide, f"  {i+1}.  {step}",
             Inches(2.5), Inches(3.98) + i * Inches(0.48),
             Inches(9.0), Inches(0.44),
             font_size=14, color=WHITE)

footer(slide, "AI Ideathon 2026  |  Stewart Title Guaranty Company  |  Live at aiideathon.vercel.app")


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Key Features Showcase
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(slide, "Key Features", "What makes this a governed enterprise system — not just a chatbot")

cards = [
    ("📎", "Cited Answers",        "Every factual claim tagged [1][2] with numbered citations. Click to expand the exact source document passage."),
    ("🎯", "Confidence Scoring",   "FULL SUPPORT — answer well covered.\nPARTIAL SUPPORT — incomplete coverage.\nNO SUPPORT — not in approved docs."),
    ("💬", "Multi-Turn Memory",    "Maintains last 3 conversation turns as context. Follow-up questions understood naturally without re-explaining."),
    ("💡", "Follow-Up Suggestions","AI proactively suggests 3 relevant next questions after every answer — keeps exploration flowing."),
    ("📊", "Session Analytics",    "Real-time dashboard: confidence breakdown, satisfaction rate from feedback, most referenced documents."),
    ("🤖", "Slack Bot /ask",       "Slash command works in any channel. Instant ACK in <1 sec. Full cited answer posted ~10 seconds later."),
]

cw = Inches(4.1)
ch = Inches(2.38)
for i, (emoji, title, body) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = Inches(0.22) + col * (cw + Inches(0.25))
    y = Inches(1.45) + row * (ch + Inches(0.22))
    icon_card(slide, emoji, title, body, x, y, cw, ch)

footer(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(slide, "System Architecture", "100% Google Cloud Platform — no data leaves your GCP project")

# Layer diagram
layers = [
    ("User Interfaces", "React Web App (Vercel)  +  Slack Bot (/ask command)", WHITE, TEAL),
    ("API Layer",       "FastAPI Backend  —  Google Cloud Run (serverless, auto-scaling)", TEAL_LIGHT, TEAL),
    ("AI Search",       "Google Discovery Engine (Vertex AI)  —  Semantic document search & retrieval", RGBColor(0xE8, 0xF4, 0xF8), RGBColor(0x00, 0x7A, 0x99)),
    ("AI Generation",   "Vertex AI  —  Gemini 2.5 Flash  —  Grounded answer synthesis + follow-ups", RGBColor(0xFD, 0xF0, 0xF2), RED),
    ("Storage",         "Google Cloud Storage  —  Approved PDFs & DOCX  |  Embeddings: text-embedding-004", RGBColor(0xF5, 0xF0, 0xFD), RGBColor(0x55, 0x33, 0x99)),
]

ly = Inches(1.42)
lh = Inches(0.88)
gap = Inches(0.1)
for label, content, bg, accent in layers:
    add_rect(slide, Inches(0.22), ly, Inches(12.89), lh, bg)
    add_rect(slide, Inches(0.22), ly, Inches(0.07), lh, accent)
    add_text(slide, label,
             Inches(0.42), ly + Inches(0.08), Inches(2.1), Inches(0.38),
             font_size=12, bold=True, color=accent)
    add_text(slide, content,
             Inches(2.6), ly + Inches(0.22), Inches(10.3), Inches(0.55),
             font_size=13, color=GRAY)
    # Arrow down
    if label != "Storage":
        add_text(slide, "↓", Inches(6.5), ly + lh, Inches(0.5), gap + Inches(0.02),
                 font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)
    ly += lh + gap

# GCP services box
add_rect(slide, Inches(0.22), Inches(6.18), Inches(12.89), Inches(0.72), TEAL_DARK)
services = ["Cloud Run", "Vertex AI Gemini 2.5 Flash", "Discovery Engine", "Cloud Storage", "Cloud Build", "Artifact Registry", "Cloud IAM"]
add_text(slide, "GCP Services:  " + "   ·   ".join(services),
         Inches(0.42), Inches(6.28), Inches(12.5), Inches(0.52),
         font_size=11.5, color=WHITE)

footer(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — RAG Pipeline
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(slide, "How Answers Are Generated", "The RAG Pipeline — step by step")

pipeline = [
    ("💬", "User Question",            "Employee types or speaks a natural language question"),
    ("🔍", "Semantic Search",          "Discovery Engine searches indexed approved documents — finds top 3 relevant passages by meaning"),
    ("📋", "Context Block Built",      "[1] Source: Cloud Infrastructure — 'The Quarter Rack has 1440 GB RAM...'\n[2] Source: App Service Guide — '...'"),
    ("🤖", "Gemini Prompted",          "\"Answer ONLY using the numbered sources below. Use citation format [1]. Do not use outside knowledge.\""),
    ("✅", "Answer + Confidence",       "Cited answer generated · Citations validated · FULL / PARTIAL / NO SUPPORT calculated · Follow-ups suggested"),
]

bw = Inches(2.35)
bh = Inches(3.55)
bx = Inches(0.22)
for i, (icon, title, body) in enumerate(pipeline):
    add_rect(slide, bx, Inches(1.45), bw, bh, TEAL_LIGHT if i % 2 == 0 else WHITE)
    add_rect(slide, bx, Inches(1.45), bw, Inches(0.07), TEAL if i % 2 == 0 else RED)
    add_text(slide, icon, bx + Inches(0.9), Inches(1.55), Inches(0.6), Inches(0.55), font_size=24, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx + Inches(0.08), Inches(2.18), bw - Inches(0.15), Inches(0.5),
             font_size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txBox = slide.shapes.add_textbox(bx + Inches(0.1), Inches(2.7), bw - Inches(0.2), Inches(2.2))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY
    run.font.name = "Calibri"
    if i < 4:
        add_text(slide, "→", bx + bw + Inches(0.01), Inches(2.9), Inches(0.18), Inches(0.5),
                 font_size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    bx += bw + Inches(0.19)

# Example response
add_rect(slide, Inches(0.22), Inches(5.1), Inches(12.89), Inches(1.22), TEAL_DARK)
add_text(slide, "Example Response:",
         Inches(0.42), Inches(5.2), Inches(2.5), Inches(0.35),
         font_size=12, bold=True, color=WHITE)
add_text(slide,
         '{ answer: "The Quarter Rack has 1440 GB RAM [1].",  confidence: "FULL_SUPPORT",  '
         'sources: [{id:1, title:"Cloud Infrastructure", snippet:"..."}],  '
         'follow_up_questions: ["What about the Half Rack?", ...] }',
         Inches(0.42), Inches(5.55), Inches(12.5), Inches(0.65),
         font_size=11, color=RGBColor(0xCC, 0xE8, 0xF0), italic=True)

footer(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — vs. M365 Copilot
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(slide, "vs. Microsoft 365 Copilot Agent", "An honest comparison — they solve similar problems, but not in the same way")

# Context note
add_rect(slide, Inches(0.35), Inches(1.38), Inches(12.63), Inches(0.42), RGBColor(0xFF, 0xF3, 0xCD))
add_text(slide, "⚠️  M365 Copilot Agent now answers from SharePoint docs and says 'not found' when missing — but still shows: \"AI-generated content may be incorrect\" on every response.",
         Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.38),
         font_size=11, color=RGBColor(0x7A, 0x4F, 0x00), bold=False)

# Table
headers = ["Feature", "M365 Copilot Agent", "Our System"]
rows = [
    ["Answers from internal docs",   "✓  SharePoint documents",               "✓  GCP-approved documents only"],
    ["Says 'not found' when missing", "✓  Yes",                               "✓  Yes — with NO SUPPORT badge"],
    ["Accuracy disclaimer",           "✗  'May be incorrect' on every answer", "✓  Explicit confidence score per answer"],
    ["Citation level",                "✗  Document-level chips",               "✓  Sentence-level numbered [1][2]"],
    ["Confidence scoring",            "✗  None",                               "✓  FULL / PARTIAL / NO SUPPORT"],
    ["Session analytics",             "✗  None",                               "✓  Built-in dashboard"],
    ["Access channels",               "✗  Teams only",                         "✓  Web app + Slack — no license needed"],
    ["Cost (100 users/year)",         "✗  ~$36,000 (Copilot license)",         "✓  ~$200 (GCP pay-per-use)"],
]

col_widths = [Inches(3.1), Inches(4.35), Inches(4.35)]
row_h = Inches(0.46)
tx = Inches(0.35)
ty = Inches(1.88)

# Header row
add_rect(slide, tx, ty, sum(col_widths) + Inches(0.3), row_h, TEAL)
cx = tx
for j, h in enumerate(headers):
    add_text(slide, h, cx + Inches(0.1), ty + Inches(0.07),
             col_widths[j] - Inches(0.1), row_h - Inches(0.1),
             font_size=12, bold=True, color=WHITE)
    cx += col_widths[j]

RED_LIGHT = RGBColor(0xFD, 0xF0, 0xF2)
for i, row in enumerate(rows):
    ry = ty + (i + 1) * row_h
    bg = TEAL_LIGHT if i % 2 == 0 else WHITE
    add_rect(slide, tx, ry, sum(col_widths) + Inches(0.3), row_h, bg)
    cx = tx
    for j, cell in enumerate(row):
        if j == 1:
            c = RED if cell.startswith("✗") else RGBColor(0x16, 0x7A, 0x50)
        elif j == 2:
            c = TEAL
        else:
            c = GRAY
        add_text(slide, cell, cx + Inches(0.1), ry + Inches(0.07),
                 col_widths[j] - Inches(0.1), row_h - Inches(0.1),
                 font_size=11.5, color=c, bold=(j == 2))
        cx += col_widths[j]

# Bottom pitch
add_rect(slide, Inches(0.35), Inches(6.0), Inches(12.63), Inches(0.72), RED)
add_text(slide,
         "\"Microsoft validates the problem. Our system solves it better — with confidence scoring, "
         "sentence-level citations, and no disclaimer saying the answer might be wrong.\"",
         Inches(0.55), Inches(6.08), Inches(12.2), Inches(0.62),
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

footer(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Business Value & Impact
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(slide, "Business Value & Impact", "A production system solving a real Stewart problem — delivered in the Ideathon timeframe")

metrics = [
    ("⚡", "< 10 seconds",  "From question to\ncited answer"),
    ("🔒", "Zero",           "Hallucination risk\nby design"),
    ("📎", "100%",           "Answers traceable\nto source"),
    ("💰", "~$200/year",     "vs $36,000/year\nfor M365 Copilot"),
]

mw = Inches(2.9)
mh = Inches(2.2)
mx = Inches(0.4)
for emoji, stat, label in metrics:
    add_rect(slide, mx, Inches(1.45), mw, mh, TEAL)
    add_text(slide, emoji, mx + Inches(1.1), Inches(1.55), Inches(0.7), Inches(0.65), font_size=26, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, stat, mx + Inches(0.1), Inches(2.25), mw - Inches(0.2), Inches(0.7),
             font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, mx + Inches(0.1), Inches(2.95), mw - Inches(0.2), Inches(0.6),
             font_size=12, color=RGBColor(0xCC, 0xE8, 0xF0), align=PP_ALIGN.CENTER)
    mx += mw + Inches(0.2)

# Value props
vx = Inches(0.22)
vy = Inches(3.85)
value_items = [
    ("🏢", "Enterprise Ready",      "Deployed on GCP. No data leaves your environment. IAM-controlled access. Compliance-friendly by architecture."),
    ("📈", "Continuous Improvement", "Analytics dashboard shows which questions fail. Knowledge managers see exactly which document gaps to fill."),
    ("🔌", "Multi-Channel Access",   "Web app for desktop users. Slack bot for teams. Same governed engine behind both — consistent answers everywhere."),
    ("🚀", "Expandable",             "Add any PDF, DOCX, or text file. Create department-specific data stores: HR, Legal, Compliance, Infrastructure."),
]
vw = Inches(3.1)
for i, (icon, title, body) in enumerate(value_items):
    vxi = vx + i * (vw + Inches(0.2))
    add_rect(slide, vxi, vy, vw, Inches(2.45), LIGHT_GRAY)
    add_rect(slide, vxi, vy, vw, Inches(0.06), RED)
    add_text(slide, icon + "  " + title, vxi + Inches(0.12), vy + Inches(0.12),
             vw - Inches(0.2), Inches(0.42),
             font_size=12, bold=True, color=TEAL)
    txBox = slide.shapes.add_textbox(vxi + Inches(0.12), vy + Inches(0.58),
                                     vw - Inches(0.22), Inches(1.75))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(11.5)
    run.font.color.rgb = GRAY
    run.font.name = "Calibri"

footer(slide)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Thank You / Closing
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, TEAL_DARK)
add_rect(slide, 0, 0, Inches(0.55), SLIDE_H, RED)
add_rect(slide, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), RED)

add_text(slide, "Thank You",
         Inches(0.8), Inches(0.75), Inches(11.7), Inches(1.1),
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(3.5), Inches(1.85), Inches(6.33), Inches(0.05), RED)

# Summary bullets
summary = [
    "✓  Production-deployed RAG system — live at aiideathon.vercel.app",
    "✓  Every answer cited, confidence-rated, and traceable to source",
    "✓  Built entirely on Google Cloud Platform — no data leaves GCP",
    "✓  Web app + Slack bot — accessible where employees already work",
    "✓  Analytics dashboard — continuous improvement built-in",
    "✓  Governance first — only approved documents can ever appear in answers",
]
for i, s in enumerate(summary):
    add_text(slide, s,
             Inches(1.5), Inches(2.1) + i * Inches(0.52),
             Inches(10.5), Inches(0.48),
             font_size=14, color=WHITE)

add_rect(slide, Inches(0.8), Inches(5.35), Inches(11.73), Inches(0.06), RED)

add_text(slide,
         "\"The question isn't whether AI can help employees find answers faster.\n"
         "The question is whether we can do it safely, traceably, and governed.\n"
         "We can. And we did.\"",
         Inches(0.8), Inches(5.52), Inches(11.73), Inches(1.0),
         font_size=14, color=RGBColor(0xCC, 0xE8, 0xF0),
         align=PP_ALIGN.CENTER, italic=True)

# Contact row
add_text(slide, "🌐  aiideathon.vercel.app   |   Meet P.   |   AI Ideathon 2026 — Stewart Title Guaranty Company",
         Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.48),
         font_size=12, color=RGBColor(0xAA, 0xCC, 0xD8), align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────
output_path = "/workspaces/AI_Ideathon/AI_Ideathon_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
